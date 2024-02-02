from pptx import Presentation
from transformers import pipeline

# Hugging Face의 텍스트 생성 파이프라인 설정
generator = pipeline('text-generation', model='gpt2')

def generate_content(topic):
    # 주제에 대한 내용 생성
    prompt = f"Explain about {topic}:"
    results = generator(prompt, max_length=100, num_return_sequences=1)
    return results[0]['generated_text']

def create_presentation(topic, content):
    # 프레젠테이션 객체 생성
    prs = Presentation()

    # 타이틀 슬라이드
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.placeholders[0]
    subtitle = slide.placeholders[1]
    title.text = "주제: " + topic
    subtitle.text = "자동 생성된 프레젠테이션"

    # 내용 슬라이드
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.placeholders[0]
    content_box = slide.placeholders[1]
    title.text = topic
    content_box.text = content

    # 프레젠테이션 저장
    prs.save(topic + '__presentation.pptx')

# 사용자 입력
topic = input("프레젠테이션 주제를 입력하세요: ")
content = generate_content(topic)

# 프레젠테이션 생성
create_presentation(topic, content)
