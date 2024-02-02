from pptx import Presentation
from pptx.util import Pt

def create_presentation(title, subtitles):
    # 프레젠테이션 객체 생성
    prs = Presentation()

    # 타이틀 슬라이드 추가
    slide_layout = prs.slide_layouts[0]  # 타이틀 슬라이드 레이아웃
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.placeholders[0]
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    subtitle_placeholder.text = "Created using Python"

    # 내용 슬라이드 추가
    for subtitle in subtitles:
        slide_layout = prs.slide_layouts[1]  # 내용 슬라이드 레이아웃
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.placeholders[0]
        content_placeholder = slide.placeholders[1]
        title_placeholder.text = subtitle
        content_placeholder.text = "자세한 내용"

    # 프레젠테이션 저장
    prs.save('presentation.pptx')

# 사용자 입력
presentation_title = "프레젠테이션 제목"
subtitles = ["첫 번째 슬라이드", "두 번째 슬라이드", "세 번째 슬라이드"]

# 프레젠테이션 생성
create_presentation(presentation_title, subtitles)
