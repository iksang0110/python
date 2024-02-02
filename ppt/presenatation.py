import openai
import os
from pptx import Presentation

# OpenAI API 키 설정
openai.api_key = os.getenv("OPENAI_API_KEY")

def generate_content(topic):
    # OpenAI GPT-3를 사용하여 주제에 대한 내용 생성
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=f"{topic}에 대한 설명을 작성해주세요.",
        max_tokens=100
    )
    return response.choices[0].text.strip()

def create_presentation(topic, content):
    # 프레젠테이션 객체 생성
    prs = Presentation()

    # 타이틀 슬라이드
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.placeholders[0]
    subtitle = slide.placeholders[1]
    title.text = "주제: " + topic
    subtitle.text = "자동 생성된 프레젠테이션"

    # 내용 슬라이드
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.placeholders[0]
    content_box = slide.placeholders[1]
    title.text = topic
    content_box.text = content

    # 프레젠테이션 저장
    prs.save(topic + '_presentation.pptx')

# 사용자 입력
topic = input("프레젠테이션 주제를 입력하세요: ")
content = generate_content(topic)

# 프레젠테이션 생성
create_presentation(topic, content)
